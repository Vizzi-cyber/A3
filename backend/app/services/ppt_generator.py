"""
PPT 自动生成服务 v2
深度内容：代码示例、算法步骤、复杂度分析、语法讲解
"""
import json
import os
import asyncio
import re
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional

from ..core.logger import setup_logger

logger = setup_logger()

PPT_OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "..", "..", "static", "resources", "ppt")
os.makedirs(PPT_OUTPUT_DIR, exist_ok=True)


# ============================================================
# LLM 大纲生成
# ============================================================

async def generate_ppt_outline(
    topic: str,
    subject: str = "C语言数据结构",
    llm=None,
) -> Dict[str, Any]:
    """用 LLM 生成深度教学 PPT 大纲"""

    prompt = f"""你是一位资深的{subject}教师，需要为学生制作一份关于「{topic}」的深度学习课件。
目标学生是计算机专业本科生，课件需要兼顾概念讲解和实战代码。

请严格按照以下JSON格式输出，不要输出任何其他内容：

```json
{{
  "title": "课件主标题",
  "subtitle": "副标题（学科名+主题定位）",
  "slides": [
    {{
      "type": "title",
      "title": "标题",
      "subtitle": "副标题"
    }},
    {{
      "type": "definition",
      "title": "概念定义",
      "concept": "核心概念名称",
      "definition": "准确定义（2-3句话）",
      "key_properties": ["性质1", "性质2", "性质3"],
      "analogy": "生活类比帮助理解（一句话）"
    }},
    {{
      "type": "structure",
      "title": "数据结构/语法详解",
      "description": "结构说明",
      "fields": [
        {{"name": "字段名", "type": "类型", "desc": "作用说明"}}
      ],
      "code": "C语言结构体定义代码（完整可编译）"
    }},
    {{
      "type": "algorithm",
      "title": "算法/操作步骤",
      "description": "算法概述",
      "steps": [
        "步骤1：具体描述",
        "步骤2：具体描述",
        "步骤3：具体描述"
      ],
      "code": "完整的C语言实现代码（带注释，可编译运行）",
      "complexity": {{
        "time": "O(?)",
        "space": "O(?)",
        "explanation": "为什么是这个复杂度"
      }}
    }},
    {{
      "type": "code_walkthrough",
      "title": "代码逐行解析",
      "code": "关键代码片段",
      "explanations": [
        {{"lines": "第1-3行", "explain": "这段代码的作用"}},
        {{"lines": "第4-6行", "explain": "这段代码的作用"}}
      ]
    }},
    {{
      "type": "comparison",
      "title": "对比分析",
      "headers": ["对比项", "方案A", "方案B"],
      "rows": [
        ["时间复杂度", "O(n)", "O(log n)"],
        ["空间复杂度", "O(1)", "O(n)"],
        ["适用场景", "...", "..."]
      ]
    }},
    {{
      "type": "knowledge_map",
      "title": "知识结构图",
      "items": ["中心概念", "子概念1", "子概念2", "子概念3"],
      "relations": [["中心概念", "子概念1"], ["中心概念", "子概念2"]]
    }},
    {{
      "type": "mistakes",
      "title": "常见易错点",
      "items": [
        {{"point": "错误描述", "correct": "正确做法", "code_example": "错误/正确代码对比（可选）"}}
      ]
    }},
    {{
      "type": "quiz",
      "title": "练习题",
      "questions": [
        {{"q": "题目", "options": ["A. ...", "B. ...", "C. ...", "D. ..."], "answer": "A", "explain": "详细解析"}}
      ]
    }},
    {{
      "type": "summary",
      "title": "总结",
      "points": ["核心要点1", "核心要点2", "核心要点3"],
      "next_topic": "下一个建议学习的主题"
    }}
  ]
}}
```

要求（非常重要）：
1. 生成10-15页幻灯片，内容要有深度，不能泛泛而谈
2. 必须包含：定义页（含性质和类比）、数据结构/语法详解页（含C语言结构体代码）、至少2页算法步骤（含完整C代码和复杂度分析）、代码逐行解析页、对比分析页、知识结构图、易错点（含代码示例）、2-3道练习题（含详细解析）、总结页
3. 所有代码必须是完整的、可编译的C语言代码，不能用伪代码
4. 复杂度分析要解释原因，不能只写O(n)
5. 易错点要具体，最好附带错误代码和正确代码对比
6. 练习题难度要有梯度：1道基础题、1道应用题、1道综合题
7. 内容必须准确，符合{subject}的标准教材定义
"""

    if llm is None:
        return _default_outline(topic, subject)

    try:
        result = await llm.ainvoke(
            messages=[{"role": "user", "content": prompt}],
            temperature=0.6,
            max_tokens=4000,
        )
        text = result if isinstance(result, str) else getattr(result, "content", str(result))
        outline = _parse_outline_json(text)
        if outline and len(outline.get("slides", [])) >= 6:
            return outline
        logger.warning("LLM PPT 大纲解析失败或页数不足，使用默认大纲")
        return _default_outline(topic, subject)
    except Exception as e:
        logger.error(f"生成 PPT 大纲失败: {e}")
        return _default_outline(topic, subject)


def _parse_outline_json(text: str) -> Optional[Dict[str, Any]]:
    """从 LLM 输出中解析 JSON"""
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        pass
    match = re.search(r"```(?:json)?\s*([\s\S]*?)```", text, re.IGNORECASE)
    if match:
        try:
            return json.loads(match.group(1))
        except json.JSONDecodeError:
            pass
    # 尝试找到第一个 { 和最后一个 }
    start = text.find("{")
    end = text.rfind("}")
    if start != -1 and end != -1 and end > start:
        try:
            return json.loads(text[start:end + 1])
        except json.JSONDecodeError:
            pass
    return None


# ============================================================
# 数据结构主题的深度默认大纲
# ============================================================

TOPIC_OUTLINES: Dict[str, Dict[str, Any]] = {
    "二叉树": {
        "title": "二叉树",
        "subtitle": "树形结构的核心 · 递归思想的基石",
        "slides": [
            {"type": "title", "title": "二叉树详解", "subtitle": "数据结构 · 从概念到实现"},
            {"type": "definition", "title": "什么是二叉树", "concept": "二叉树 (Binary Tree)",
             "definition": "二叉树是每个节点最多有两个子树的树结构，通常称为左子树和右子树。二叉树是递归定义的：它要么为空，要么由一个根节点和两棵互不相交的二叉树组成。",
             "key_properties": ["每个节点最多2个子节点", "第i层最多2^(i-1)个节点", "深度为k的二叉树最多2^k-1个节点", "叶子节点数 = 度为2的节点数 + 1"],
             "analogy": "就像家族族谱，每个人最多有两个孩子，层层向下传承"},
            {"type": "structure", "title": "二叉树的存储结构", "description": "链式存储是最常用的方式，每个节点包含数据域和两个指针域",
             "fields": [{"name": "data", "type": "ElementType", "desc": "存储节点数据"}, {"name": "left", "type": "BiTreeNode*", "desc": "指向左子节点的指针"}, {"name": "right", "type": "BiTreeNode*", "desc": "指向右子节点的指针"}],
             "code": "typedef struct BiTreeNode {\n    int data;\n    struct BiTreeNode *left;\n    struct BiTreeNode *right;\n} BiTreeNode, *BiTree;"},
            {"type": "algorithm", "title": "前序遍历（递归）", "description": "访问顺序：根 → 左 → 右。是最基础的树遍历方式",
             "steps": ["步骤1：访问当前节点（输出/处理数据）", "步骤2：递归遍历左子树", "步骤3：递归遍历右子树", "递归出口：当前节点为NULL时返回"],
             "code": "void preOrder(BiTree T) {\n    if (T == NULL) return;\n    printf(\"%d \", T->data);  // 访问根\n    preOrder(T->left);        // 遍历左子树\n    preOrder(T->right);       // 遍历右子树\n}",
             "complexity": {"time": "O(n)", "space": "O(h)，h为树高", "explanation": "每个节点恰好访问一次，所以时间O(n)；空间复杂度取决于递归栈深度，最坏情况（链状树）为O(n)，平衡树为O(logn)"}},
            {"type": "algorithm", "title": "层序遍历（队列实现）", "description": "按层次从上到下、从左到右访问每个节点，需要借助队列",
             "steps": ["步骤1：根节点入队", "步骤2：队头节点出队并访问", "步骤3：将出队节点的左、右子节点依次入队", "步骤4：重复步骤2-3直到队列为空"],
             "code": "void levelOrder(BiTree T) {\n    if (T == NULL) return;\n    Queue Q;\n    InitQueue(&Q);\n    EnQueue(&Q, T);\n    while (!IsEmpty(Q)) {\n        BiTreeNode *node = DeQueue(&Q);\n        printf(\"%d \", node->data);\n        if (node->left) EnQueue(&Q, node->left);\n        if (node->right) EnQueue(&Q, node->right);\n    }\n}",
             "complexity": {"time": "O(n)", "space": "O(n)", "explanation": "每个节点入队出队各一次，时间O(n)；队列最多存储一层的节点数，最坏情况（完全二叉树最后一层）约n/2个，所以O(n)"}},
            {"type": "code_walkthrough", "title": "递归建树代码解析",
             "code": "BiTree createTree() {\n    int val;\n    scanf(\"%d\", &val);\n    if (val == -1) return NULL;  // -1表示空节点\n    BiTree T = (BiTree)malloc(sizeof(BiTreeNode));\n    T->data = val;\n    T->left = createTree();     // 递归创建左子树\n    T->right = createTree();    // 递归创建右子树\n    return T;\n}",
             "explanations": [{"lines": "第3-4行", "explain": "读取输入值，-1作为哨兵值表示该位置为空节点，这是递归出口"}, {"lines": "第5-6行", "explain": "动态分配节点内存，设置数据域。malloc返回void*，自动转换为BiTree类型"}, {"lines": "第7-8行", "explain": "核心递归：分别递归创建左右子树，体现了二叉树的递归定义"}]},
            {"type": "comparison", "title": "三种遍历方式对比",
             "headers": ["遍历方式", "访问顺序", "典型应用", "实现方式"],
             "rows": [["前序遍历", "根→左→右", "复制二叉树、前缀表达式", "递归/栈"], ["中序遍历", "左→根→右", "BST排序输出、中缀表达式", "递归/栈"], ["后序遍历", "左→右→根", "释放内存、后缀表达式、计算目录大小", "递归/栈"], ["层序遍历", "逐层从左到右", "求树的宽度、最短路径", "队列"]]},
            {"type": "knowledge_map", "title": "二叉树知识结构", "items": ["二叉树", "存储结构", "遍历算法", "特殊二叉树", "应用场景"],
             "relations": [["二叉树", "存储结构"], ["二叉树", "遍历算法"], ["二叉树", "特殊二叉树"], ["遍历算法", "应用场景"], ["存储结构", "遍历算法"]]},
            {"type": "mistakes", "title": "常见易错点",
             "items": [
                 {"point": "递归遍历忘记写递归出口", "correct": "必须先判断 if(T == NULL) return; 否则会段错误", "code_example": "❌ void preOrder(BiTree T) {\n       printf(\"%d\", T->data);\n       preOrder(T->left);\n   }\n✅ void preOrder(BiTree T) {\n       if (T == NULL) return;\n       printf(\"%d\", T->data);\n       preOrder(T->left);\n   }"},
                 {"point": "混淆前中后序遍历的访问位置", "correct": "记住：前序=根在前，中序=根在中，后序=根在后。visit()相对于左右递归的位置决定遍历类型"},
                 {"point": "层序遍历误用栈而非队列", "correct": "层序遍历必须用队列（FIFO），用栈会变成另一种遍历顺序"}
             ]},
            {"type": "quiz", "title": "练习题",
             "questions": [
                 {"q": "一棵完全二叉树有100个节点，其叶子节点数为多少？", "options": ["A. 49", "B. 50", "C. 51", "D. 52"], "answer": "B", "explain": "完全二叉树中，n0 = n/2（向下取整）。100/2 = 50。也可以用公式：n0 = n2 + 1，而n = n0 + n1 + n2，完全二叉树n1只能为0或1"},
                 {"q": "前序遍历序列为ABDECFG，中序遍历序列为DBEAFCG，后序遍历序列是什么？", "options": ["A. DEBFGCA", "B. DEBFACG", "C. DEBFCGA", "D. DEBFCA G"], "answer": "A", "explain": "前序第一个A是根，在中序中找到A的位置，左边DBE是左子树，右边FCG是右子树。递归处理：左子树前序BDE→中序DBE，右子树CFG→中序FCG。后序：左→右→根 = DEB FGC A"},
                 {"q": "以下关于二叉树的说法，错误的是？", "options": ["A. 二叉树可以为空", "B. 二叉树的度一定不超过2", "C. 二叉树是有序树", "D. 满二叉树就是完全二叉树"], "answer": "D", "explain": "满二叉树是每一层都达到最大节点数的二叉树，完全二叉树只有最后一层可以不满且从左到右连续。满二叉树是完全二叉树的特例，但反过来不成立。D选项说'就是'意味着二者等同，这是错误的"}
             ]},
            {"type": "summary", "title": "总结", "points": ["二叉树是递归定义的树形结构，最多两个子节点", "四种遍历方式各有特点：前序适合复制，中序适合排序，后序适合释放，层序适合求宽度", "递归是二叉树操作的核心思想，务必掌握递归出口", "完全二叉树和满二叉树是重要的特殊形态，有数学性质"], "next_topic": "二叉搜索树 (BST) 或 哈夫曼树"},
        ],
    },
    "链表": {
        "title": "链表",
        "subtitle": "动态存储的基石 · 指针操作的核心训练",
        "slides": [
            {"type": "title", "title": "链表详解", "subtitle": "数据结构 · 指针与动态内存"},
            {"type": "definition", "title": "什么是链表", "concept": "链表 (Linked List)",
             "definition": "链表是一种物理存储单元上非连续、非顺序的存储结构，数据元素的逻辑顺序通过链表中的指针链接次序实现。每个节点包含数据域和指针域。",
             "key_properties": ["动态分配内存，不需要预先知道数据量", "插入删除只需修改指针，时间O(1)", "不支持随机访问，查找需要O(n)", "没有空间浪费（对比数组的预分配）"],
             "analogy": "就像寻宝游戏，每个线索告诉你下一个线索在哪里，必须顺着走完"},
            {"type": "structure", "title": "单链表节点定义", "description": "每个节点包含数据和指向下一个节点的指针",
             "fields": [{"name": "data", "type": "int", "desc": "存储节点数据"}, {"name": "next", "type": "struct Node*", "desc": "指向下一个节点的指针，最后一个节点指向NULL"}],
             "code": "typedef struct Node {\n    int data;\n    struct Node *next;\n} Node, *LinkList;"},
            {"type": "algorithm", "title": "头插法建表", "description": "每次将新节点插入到链表头部，最终得到的顺序与输入顺序相反",
             "steps": ["步骤1：创建头结点（不存数据，next指向NULL）", "步骤2：读取新数据，创建新节点", "步骤3：新节点的next指向原第一个节点", "步骤4：头结点的next指向新节点", "步骤5：重复步骤2-4直到所有数据插入完毕"],
             "code": "LinkList createListHead(int arr[], int n) {\n    LinkList L = (LinkList)malloc(sizeof(Node));\n    L->next = NULL;  // 创建头结点\n    for (int i = 0; i < n; i++) {\n        Node *s = (Node*)malloc(sizeof(Node));\n        s->data = arr[i];\n        s->next = L->next;  // 新节点指向原首节点\n        L->next = s;        // 头结点指向新节点\n    }\n    return L;\n}",
             "complexity": {"time": "O(n)", "space": "O(n)", "explanation": "n个元素各插入一次，每次插入是O(1)，总共O(n)；需要n个节点的存储空间"}},
            {"type": "algorithm", "title": "在第i个位置插入节点", "description": "核心操作：先找到第i-1个节点，再修改指针完成插入",
             "steps": ["步骤1：从头结点开始，用指针p遍历到第i-1个节点", "步骤2：检查位置是否合法（p不能为NULL）", "步骤3：创建新节点s，设置数据域", "步骤4：s->next = p->next（新节点指向原第i个节点）", "步骤5：p->next = s（第i-1个节点指向新节点）"],
             "code": "Status insert(LinkList L, int i, int e) {\n    Node *p = L;\n    int j = 0;\n    while (p && j < i - 1) {  // 找到第i-1个节点\n        p = p->next;\n        j++;\n    }\n    if (!p || j > i - 1) return ERROR;  // 位置非法\n    Node *s = (Node*)malloc(sizeof(Node));\n    s->data = e;\n    s->next = p->next;  // ① 新节点指向后继\n    p->next = s;        // ② 前驱指向新节点\n    return OK;\n}",
             "complexity": {"time": "O(n)", "space": "O(1)", "explanation": "查找第i-1个节点需要遍历，最坏O(n)；插入操作本身只修改两个指针，O(1)"}},
            {"type": "code_walkthrough", "title": "删除节点代码解析",
             "code": "Status delete(LinkList L, int i, int *e) {\n    Node *p = L;\n    int j = 0;\n    while (p->next && j < i - 1) {\n        p = p->next;\n        j++;\n    }\n    if (!(p->next) || j > i - 1) return ERROR;\n    Node *q = p->next;     // q指向待删除节点\n    *e = q->data;           // 保存被删节点的值\n    p->next = q->next;     // 前驱跳过被删节点\n    free(q);                // 释放内存\n    return OK;\n}",
             "explanations": [{"lines": "第4行", "explain": "循环条件用p->next而非p，因为我们需要确保待删除节点存在"}, {"lines": "第9-10行", "explain": "先用q保存待删除节点的地址，再取出数据"}, {"lines": "第11行", "explain": "核心操作：让前驱节点直接指向后继节点，跳过被删节点"}, {"lines": "第12行", "explain": "必须free释放内存，否则会造成内存泄漏"}]},
            {"type": "comparison", "title": "数组 vs 链表",
             "headers": ["对比项", "数组", "链表"],
             "rows": [["内存分配", "连续空间，编译时确定", "离散空间，运行时动态分配"], ["随机访问", "O(1) 支持", "O(n) 不支持"], ["插入删除", "O(n) 需移动元素", "O(1) 只改指针"], ["空间利用", "可能浪费（预分配）", "按需分配，无浪费"], ["适用场景", "频繁查找、已知数据量", "频繁增删、数据量不确定"]]},
            {"type": "mistakes", "title": "常见易错点",
             "items": [
                 {"point": "插入操作指针顺序颠倒", "correct": "必须先让新节点指向后继，再让前驱指向新节点。反过来会丢失后继节点地址", "code_example": "❌ p->next = s; s->next = p->next; // s->next指向了自己！\n✅ s->next = p->next; p->next = s;"},
                 {"point": "删除节点后忘记free", "correct": "free(q)释放内存是必须的，否则内存泄漏"},
                 {"point": "遍历时丢失头结点引用", "correct": "用工作指针p遍历，不要移动头指针L，否则链表就断了"}
             ]},
            {"type": "quiz", "title": "练习题",
             "questions": [
                 {"q": "单链表中，增加头结点的目的是什么？", "options": ["A. 使遍历更方便", "B. 使空表和非空表的处理统一", "C. 节省存储空间", "D. 提高查找效率"], "answer": "B", "explain": "头结点不存储数据，它的存在使得第一个数据节点的插入删除操作与其他节点一致，不需要特殊处理空表的情况"},
                 {"q": "在一个长度为n的单链表中，删除已知指针p所指节点的时间复杂度是？", "options": ["A. O(1)", "B. O(logn)", "C. O(n)", "D. O(n²)"], "answer": "A", "explain": "如果已知p指向待删除节点，并且我们用'将后继节点值复制到p再删除后继'的技巧，可以在O(1)完成。但如果是'给定值删除'则需要O(n)查找"},
                 {"q": "以下代码执行后，链表的状态是什么？\nNode *s = p->next; p->next = s->next; free(s);", "options": ["A. 删除了p节点", "B. 删除了p的后继节点", "C. 链表断裂", "D. 删除了p的前驱节点"], "answer": "B", "explain": "s保存了p的下一个节点，p->next = s->next让p跳过s指向后面的节点，最后free(s)释放s。所以删除的是p的直接后继"}
             ]},
            {"type": "summary", "title": "总结", "points": ["链表通过指针实现动态存储，插入删除效率高", "头插法/尾插法是建表的两种基本方式", "插入删除操作的核心是正确修改指针顺序", "头结点简化了边界条件的处理"], "next_topic": "双向链表 或 循环链表"},
        ],
    },
    "栈": {
        "title": "栈",
        "subtitle": "后进先出的线性结构 · 递归与回溯的基础",
        "slides": [
            {"type": "title", "title": "栈详解", "subtitle": "数据结构 · LIFO 原理与应用"},
            {"type": "definition", "title": "什么是栈", "concept": "栈 (Stack)",
             "definition": "栈是一种仅允许在一端（栈顶）进行插入和删除操作的线性表。后进先出（LIFO）是其核心特性。",
             "key_properties": ["后进先出 LIFO", "只在栈顶操作", "push 入栈 / pop 出栈", "应用：递归、表达式求值、括号匹配"],
             "analogy": "像弹夹：最后压入的子弹最先射出"},
            {"type": "structure", "title": "栈的存储结构", "description": "顺序栈和链栈是两种基本实现方式",
             "fields": [{"name": "data", "type": "int[]", "desc": "存储栈元素的数组"}, {"name": "top", "type": "int", "desc": "栈顶指针，指向栈顶元素位置"}, {"name": "maxSize", "type": "int", "desc": "栈的最大容量"}],
             "code": "#define MAXSIZE 100\ntypedef struct {\n    int data[MAXSIZE];\n    int top;  // 栈顶指针\n} SqStack;"},
            {"type": "algorithm", "title": "入栈与出栈操作", "description": "栈的核心操作：push 和 pop",
             "steps": ["push: 检查栈是否满 → top++ → 存入元素", "pop: 检查栈是否空 → 取出元素 → top--", "peek: 返回栈顶元素但不修改 top", "isEmpty: 判断 top 是否为 -1"],
             "code": "Status Push(SqStack *S, int e) {\n    if (S->top == MAXSIZE - 1) return ERROR;\n    S->top++;\n    S->data[S->top] = e;\n    return OK;\n}\n\nStatus Pop(SqStack *S, int *e) {\n    if (S->top == -1) return ERROR;\n    *e = S->data[S->top];\n    S->top--;\n    return OK;\n}",
             "complexity": {"time": "O(1)", "space": "O(n)", "explanation": "入栈出栈只修改栈顶指针，O(1)；存储需要O(n)空间"}},
            {"type": "algorithm", "title": "括号匹配", "description": "栈的经典应用：检验表达式中括号是否匹配",
             "steps": ["遍历表达式每个字符", "遇到左括号：入栈", "遇到右括号：检查栈顶是否匹配的左括号", "遍历结束：栈为空则匹配成功"],
             "code": "bool bracketMatch(char *expr) {\n    Stack S;\n    InitStack(&S);\n    for (int i = 0; expr[i]; i++) {\n        if (expr[i] == '(' || expr[i] == '[' || expr[i] == '{')\n            Push(&S, expr[i]);\n        else if (expr[i] == ')' || expr[i] == ']' || expr[i] == '}') {\n            if (IsEmpty(S)) return false;\n            char top;\n            Pop(&S, &top);\n            if (expr[i] == ')' && top != '(') return false;\n            if (expr[i] == ']' && top != '[') return false;\n            if (expr[i] == '}' && top != '{') return false;\n        }\n    }\n    return IsEmpty(S);\n}",
             "complexity": {"time": "O(n)", "space": "O(n)", "explanation": "遍历一次表达式O(n)；最坏情况所有字符都是左括号，栈空间O(n)"}},
            {"type": "comparison", "title": "顺序栈 vs 链栈",
             "headers": ["对比项", "顺序栈", "链栈"],
             "rows": [["存储方式", "连续数组", "离散链表"], ["容量", "固定 maxSize", "动态扩展"], ["内存", "可能浪费", "按需分配"], ["实现难度", "简单", "稍复杂"], ["适用场景", "已知最大深度", "深度不确定"]]},
            {"type": "mistakes", "title": "常见易错点",
             "items": [
                 {"point": "顺序栈判断满的条件写错", "correct": "栈满条件是 top == MAXSIZE-1，不是 top == MAXSIZE", "code_example": "❌ if (S->top == MAXSIZE) // 越界\n✅ if (S->top == MAXSIZE - 1)"},
                 {"point": "出栈前未判空", "correct": "pop 前必须检查栈是否为空，否则会访问非法内存"},
                 {"point": "混淆栈顶指针的初始值", "correct": "top=-1 表示空栈（指向栈顶元素），top=0 表示指向下一个空位"}
             ]},
            {"type": "quiz", "title": "练习题",
             "questions": [
                 {"q": "若入栈序列为1,2,3,4，则不可能的出栈序列是？", "options": ["A. 4,3,2,1", "B. 1,2,3,4", "C. 1,3,2,4", "D. 4,1,2,3"], "answer": "D", "explain": "4先出栈意味着1,2,3,4都已入栈，此时栈顶到栈底为3,2,1。接下来只能按3→2→1顺序出栈，不可能先出1再出2"},
                 {"q": "用两个栈模拟队列，最少需要多少次操作完成一次入队？", "options": ["A. 1次", "B. 2次", "C. n次", "D. 不可能"], "answer": "B", "explain": "入队：push到栈A（1次）。出队：若栈B空，将A全部pop到B（摊还O(1)），再pop B。均摊后入队O(1)"},
                 {"q": "递归调用使用的栈空间最大深度取决于？", "options": ["A. 变量个数", "B. 递归层数", "C. 数据规模", "D. 返回值大小"], "answer": "B", "explain": "每层递归在系统调用栈上占用一帧，最大深度等于最深递归层数"}
             ]},
            {"type": "summary", "title": "总结", "points": ["栈是LIFO线性结构，只在栈顶操作", "顺序栈和链栈各有适用场景", "栈的经典应用：括号匹配、表达式求值、递归", "注意判空判满和栈顶指针初始值"], "next_topic": "队列"},
        ],
    },
    "队列": {
        "title": "队列",
        "subtitle": "先进先出的线性结构 · BFS与缓冲的核心",
        "slides": [
            {"type": "title", "title": "队列详解", "subtitle": "数据结构 · FIFO 原理与实现"},
            {"type": "definition", "title": "什么是队列", "concept": "队列 (Queue)",
             "definition": "队列是一种只允许在一端（队尾）插入、另一端（队头）删除的线性表。先进先出（FIFO）是其核心特性。",
             "key_properties": ["先进先出 FIFO", "队尾入队 enQueue", "队头出队 deQueue", "应用：BFS、任务调度、缓冲区"],
             "analogy": "像排队买票：先来的先走"},
            {"type": "structure", "title": "循环队列结构", "description": "用数组实现的队列，通过取模运算解决假溢出问题",
             "fields": [{"name": "data", "type": "int[]", "desc": "存储队列元素的数组"}, {"name": "front", "type": "int", "desc": "队头指针"}, {"name": "rear", "type": "int", "desc": "队尾指针"}],
             "code": "#define MAXSIZE 100\ntypedef struct {\n    int data[MAXSIZE];\n    int front;  // 队头\n    int rear;   // 队尾\n} SqQueue;"},
            {"type": "algorithm", "title": "循环队列操作", "description": "入队、出队、判空、判满",
             "steps": ["入队: rear = (rear+1) % MAXSIZE → 存入元素", "出队: 取 front 元素 → front = (front+1) % MAXSIZE", "判空: front == rear", "判满: (rear+1) % MAXSIZE == front（牺牲一个位置）"],
             "code": "Status EnQueue(SqQueue *Q, int e) {\n    if ((Q->rear+1) % MAXSIZE == Q->front) return ERROR;\n    Q->rear = (Q->rear + 1) % MAXSIZE;\n    Q->data[Q->rear] = e;\n    return OK;\n}\n\nStatus DeQueue(SqQueue *Q, int *e) {\n    if (Q->front == Q->rear) return ERROR;\n    Q->front = (Q->front + 1) % MAXSIZE;\n    *e = Q->data[Q->front];\n    return OK;\n}",
             "complexity": {"time": "O(1)", "space": "O(n)", "explanation": "入队出队都是修改指针取模运算，O(1)"}},
            {"type": "comparison", "title": "队列的三种实现",
             "headers": ["实现方式", "优点", "缺点"],
             "rows": [["顺序队列（数组）", "简单直观", "假溢出问题"], ["循环队列", "解决假溢出", "牺牲一个空间判满"], ["链式队列", "无容量限制", "每个节点额外指针开销"]]},
            {"type": "mistakes", "title": "常见易错点",
             "items": [
                 {"point": "循环队列判满条件写错", "correct": "判满是 (rear+1)%MAXSIZE==front，不是 rear==front（那是判空）", "code_example": "❌ if (Q->rear == Q->front) // 这是判空！\n✅ if ((Q->rear+1) % MAXSIZE == Q->front)"},
                 {"point": "队列元素个数计算错误", "correct": "元素个数 = (rear - front + MAXSIZE) % MAXSIZE，不能直接 rear-front"},
                 {"point": "出队后忘记更新 front", "correct": "出队必须 front = (front+1)%MAXSIZE，否则队头元素还在"}
             ]},
            {"type": "quiz", "title": "练习题",
             "questions": [
                 {"q": "循环队列中，若MAXSIZE=10，front=2，rear=5，则队列中有几个元素？", "options": ["A. 3", "B. 5", "C. 7", "D. 8"], "answer": "A", "explain": "(rear-front+MAXSIZE)%MAXSIZE = (5-2+10)%10 = 3"},
                 {"q": "BFS遍历图使用什么数据结构？", "options": ["A. 栈", "B. 队列", "C. 数组", "D. 树"], "answer": "B", "explain": "BFS按层次遍历，先访问的节点的邻居也先被处理，符合FIFO特性"},
                 {"q": "用链表实现队列，队头指针的作用是？", "options": ["A. 指向队尾", "B. 指向队头节点", "C. 指向头结点", "D. 记录队列长度"], "answer": "B", "explain": "队头指针直接指向第一个数据节点，方便出队操作"}
             ]},
            {"type": "summary", "title": "总结", "points": ["队列是FIFO线性结构", "循环队列通过取模运算解决假溢出", "判满判空条件要区分清楚", "BFS是队列最重要的应用之一"], "next_topic": "树与二叉树"},
        ],
    },
    "排序": {
        "title": "排序算法",
        "subtitle": "从基础到高级 · 时间复杂度与稳定性的权衡",
        "slides": [
            {"type": "title", "title": "排序算法详解", "subtitle": "数据结构 · 经典排序全面对比"},
            {"type": "definition", "title": "排序概述", "concept": "排序 (Sorting)",
             "definition": "排序是将一个无序序列按照关键字的大小关系排列成有序序列的过程。排序算法的评价指标包括时间复杂度、空间复杂度和稳定性。",
             "key_properties": ["稳定性：相等元素的相对顺序是否改变", "原地排序：是否需要额外空间", "比较排序：基于比较的下界为O(nlogn)", "非比较排序：基数排序、计数排序可突破下界"],
             "analogy": "像整理书架：按不同标准（书名、作者、颜色）排列"},
            {"type": "algorithm", "title": "快速排序", "description": "分治思想的经典应用，平均性能最优的排序算法",
             "steps": ["步骤1：选择基准元素（pivot）", "步骤2：分区：比pivot小的放左边，大的放右边", "步骤3：递归对左右子数组执行快排", "递归出口：子数组长度<=1"],
             "code": "void quickSort(int arr[], int low, int high) {\n    if (low >= high) return;\n    int pivot = arr[low];  // 选第一个元素为基准\n    int i = low, j = high;\n    while (i < j) {\n        while (i < j && arr[j] >= pivot) j--;\n        while (i < j && arr[i] <= pivot) i++;\n        if (i < j) swap(&arr[i], &arr[j]);\n    }\n    swap(&arr[low], &arr[i]);  // pivot归位\n    quickSort(arr, low, i - 1);\n    quickSort(arr, i + 1, high);\n}",
             "complexity": {"time": "平均O(nlogn)，最坏O(n²)", "space": "O(logn)", "explanation": "平均情况下每次分区约均分，递归深度logn；最坏情况（已排序数组）每次只减少1个元素"}},
            {"type": "algorithm", "title": "归并排序", "description": "稳定的分治排序，空间换时间的典型代表",
             "steps": ["步骤1：将数组从中间分成两半", "步骤2：递归排序左右两半", "步骤3：将两个有序子数组合并", "合并：双指针依次比较，取较小者"],
             "code": "void merge(int arr[], int temp[], int left, int mid, int right) {\n    int i = left, j = mid + 1, k = left;\n    while (i <= mid && j <= right) {\n        if (arr[i] <= arr[j])\n            temp[k++] = arr[i++];\n        else\n            temp[k++] = arr[j++];\n    }\n    while (i <= mid) temp[k++] = arr[i++];\n    while (j <= right) temp[k++] = arr[j++];\n    for (int i = left; i <= right; i++)\n        arr[i] = temp[i];\n}\n\nvoid mergeSort(int arr[], int temp[], int left, int right) {\n    if (left >= right) return;\n    int mid = (left + right) / 2;\n    mergeSort(arr, temp, left, mid);\n    mergeSort(arr, temp, mid + 1, right);\n    merge(arr, temp, left, mid, right);\n}",
             "complexity": {"time": "O(nlogn)", "space": "O(n)", "explanation": "每次均分O(logn)层，每层合并O(n)；需要O(n)辅助数组"}},
            {"type": "comparison", "title": "经典排序算法对比",
             "headers": ["算法", "平均时间", "最坏时间", "空间", "稳定性"],
             "rows": [["冒泡排序", "O(n²)", "O(n²)", "O(1)", "稳定"], ["选择排序", "O(n²)", "O(n²)", "O(1)", "不稳定"], ["插入排序", "O(n²)", "O(n²)", "O(1)", "稳定"], ["快速排序", "O(nlogn)", "O(n²)", "O(logn)", "不稳定"], ["归并排序", "O(nlogn)", "O(nlogn)", "O(n)", "稳定"], ["堆排序", "O(nlogn)", "O(nlogn)", "O(1)", "不稳定"]]},
            {"type": "mistakes", "title": "常见易错点",
             "items": [
                 {"point": "快排最坏情况的触发条件", "correct": "当数组已排序且选第一个元素为基准时，每次分区只减少1个元素，退化为O(n²)", "code_example": "避免方法：随机选基准 或 三数取中法"},
                 {"point": "归并排序的辅助空间容易忽略", "correct": "归并需要O(n)额外空间，不适合内存受限场景"},
                 {"point": "不稳定排序的例子记混", "correct": "选择排序、快排、堆排是不稳定的。例：[5a, 5b, 3] 选择排序后可能变成 [3, 5b, 5a]"}
             ]},
            {"type": "quiz", "title": "练习题",
             "questions": [
                 {"q": "对10个记录进行快排，最坏情况下的比较次数是？", "options": ["A. 45", "B. 90", "C. 100", "D. 1023"], "answer": "B", "explain": "最坏情况每次只排除1个元素：9+8+7+...+1 = 45次分区，但每次分区内的比较次数合计约n(n-1)/2 = 45*2 = 90"},
                 {"q": "以下哪个排序算法在最好情况下时间复杂度为O(n)？", "options": ["A. 快排", "B. 归并", "C. 插入排序", "D. 堆排序"], "answer": "C", "explain": "插入排序在数组已排序时只需遍历一次，O(n)。其他算法最好情况也是O(nlogn)"},
                 {"q": "需要稳定排序且空间有限时，应选择？", "options": ["A. 快排", "B. 归并", "C. 插入排序", "D. 堆排"], "answer": "C", "explain": "插入排序稳定且O(1)空间，但O(n²)时间。归并稳定但O(n)空间。需根据数据规模权衡"}
             ]},
            {"type": "summary", "title": "总结", "points": ["快排平均O(nlogn)但不稳定，归并稳定但需O(n)空间", "没有完美的排序算法，需根据场景选择", "稳定性在多关键字排序中很重要", "理解每种算法的适用场景比记住复杂度更重要"], "next_topic": "查找算法"},
        ],
    },
    "图": {
        "title": "图",
        "subtitle": "非线性结构的核心 · 关系建模的利器",
        "slides": [
            {"type": "title", "title": "图详解", "subtitle": "数据结构 · 顶点与边的艺术"},
            {"type": "definition", "title": "什么是图", "concept": "图 (Graph)",
             "definition": "图是由顶点集合和边集合组成的数据结构，用于表示多对多的关系。分为有向图和无向图，带权图和无权图。",
             "key_properties": ["顶点(Vertex)和边(Edge)", "有向图 vs 无向图", "度(Degree)：与顶点相连的边数", "路径：从一个顶点到另一个顶点的边序列"],
             "analogy": "像城市交通网络：城市是顶点，道路是边"},
            {"type": "structure", "title": "图的存储结构", "description": "邻接矩阵和邻接表是两种主要存储方式",
             "fields": [{"name": "adj matrix", "type": "int[][]", "desc": "邻接矩阵：用二维数组表示顶点间的连接"}, {"name": "adj list", "type": "Node*[]", "desc": "邻接表：每个顶点维护一个链表存储邻接点"}],
             "code": "// 邻接矩阵\ntypedef struct {\n    int vexnum, arcnum;\n    int adj[MAX][MAX];\n} MGraph;\n\n// 邻接表\ntypedef struct ArcNode {\n    int adjvex;\n    struct ArcNode *next;\n} ArcNode;\ntypedef struct {\n    int data;\n    ArcNode *first;\n} VNode;"},
            {"type": "algorithm", "title": "深度优先搜索 DFS", "description": "沿着一条路径尽可能深地搜索，回溯后探索其他路径",
             "steps": ["步骤1：访问起始顶点，标记为已访问", "步骤2：对其未访问的邻接顶点递归执行DFS", "步骤3：回溯：当没有未访问邻接点时返回", "步骤4：重复直到所有连通顶点都被访问"],
             "code": "bool visited[MAX];\n\nvoid DFS(MGraph G, int v) {\n    visited[v] = true;\n    printf(\"%d \", v);\n    for (int w = 0; w < G.vexnum; w++) {\n        if (G.adj[v][w] && !visited[w])\n            DFS(G, w);\n    }\n}",
             "complexity": {"time": "O(n²)邻接矩阵 / O(n+e)邻接表", "space": "O(n)", "explanation": "每个顶点访问一次，邻接矩阵需遍历所有可能边，邻接表只需遍历实际边"}},
            {"type": "algorithm", "title": "广度优先搜索 BFS", "description": "逐层搜索：先访问所有距离为1的顶点，再访问距离为2的...",
             "steps": ["步骤1：起始顶点入队并标记", "步骤2：队头顶点出队并访问", "步骤3：将出队顶点的所有未访问邻接入队并标记", "步骤4：重复2-3直到队列为空"],
             "code": "void BFS(MGraph G, int v) {\n    Queue Q;\n    InitQueue(&Q);\n    visited[v] = true;\n    EnQueue(&Q, v);\n    while (!IsEmpty(Q)) {\n        int u = DeQueue(&Q);\n        printf(\"%d \", u);\n        for (int w = 0; w < G.vexnum; w++) {\n            if (G.adj[u][w] && !visited[w]) {\n                visited[w] = true;\n                EnQueue(&Q, w);\n            }\n        }\n    }\n}",
             "complexity": {"time": "O(n²) / O(n+e)", "space": "O(n)", "explanation": "与DFS相同的时间复杂度；队列空间最坏O(n)"}},
            {"type": "comparison", "title": "DFS vs BFS",
             "headers": ["对比项", "DFS", "BFS"],
             "rows": [["数据结构", "栈（递归调用栈）", "队列"], ["搜索策略", "深度优先，一条路走到黑", "广度优先，逐层扩展"], ["空间", "O(h)递归深度", "O(w)最大层宽"], ["最短路径", "不保证", "无权图保证"], ["适用场景", "拓扑排序、连通分量", "最短路径、层序遍历"]]},
            {"type": "mistakes", "title": "常见易错点",
             "items": [
                 {"point": "DFS和BFS的visited标记时机", "correct": "入栈/入队时就标记，不是出栈/出队时。否则会重复入队", "code_example": "❌ 出队时标记 → 同一顶点可能多次入队\n✅ 入队时标记 → 每个顶点只入队一次"},
                 {"point": "邻接矩阵和邻接表的遍历复杂度混淆", "correct": "邻接矩阵必须检查所有n²个元素；邻接表只需检查实际边数e"},
                 {"point": "有向图和无向图的度数计算", "correct": "无向图度=边数；有向图分出度和入度"}
             ]},
            {"type": "quiz", "title": "练习题",
             "questions": [
                 {"q": "一个有n个顶点e条边的无向图，用邻接表存储，DFS的时间复杂度是？", "options": ["A. O(n)", "B. O(n²)", "C. O(n+e)", "D. O(ne)"], "answer": "C", "explain": "每个顶点访问一次O(n)，每条边检查两次（无向图）O(e)，合计O(n+e)"},
                 {"q": "BFS可以用来求解什么问题？", "options": ["A. 拓扑排序", "B. 无权图最短路径", "C. 关键路径", "D. 最小生成树"], "answer": "B", "explain": "BFS按层扩展，第一次到达某顶点的路径就是最短路径（无权图）"},
                 {"q": "以下关于图的说法，错误的是？", "options": ["A. 有向图的邻接矩阵不一定对称", "B. 无向图的邻接矩阵一定对称", "C. BFS需要队列", "D. DFS只能用递归实现"], "answer": "D", "explain": "DFS可以用显式栈非递归实现，不一定用递归"}
             ]},
            {"type": "summary", "title": "总结", "points": ["图是多对多关系的数学模型", "邻接矩阵适合稠密图，邻接表适合稀疏图", "DFS用栈/递归，BFS用队列", "BFS可求无权图最短路径"], "next_topic": "最短路径算法 (Dijkstra)"},
        ],
    },
}


def _default_outline(topic: str, subject: str) -> Dict[str, Any]:
    """根据主题匹配深度大纲，匹配不到则生成通用大纲"""
    # 尝试匹配预设大纲（支持部分匹配）
    topic_lower = topic.lower()
    for key, outline in TOPIC_OUTLINES.items():
        if key in topic_lower or topic_lower in key:
            return outline
    # 扩展匹配：排序相关、查找相关、树相关
    if any(k in topic_lower for k in ["排序", "快排", "归并", "冒泡", "选择排序", "插入排序", "堆排"]):
        return TOPIC_OUTLINES["排序"]
    if any(k in topic_lower for k in ["查找", "二分", "搜索"]):
        pass  # fall through to generic outline below
    if any(k in topic_lower for k in ["栈", "stack"]):
        return TOPIC_OUTLINES["栈"]
    if any(k in topic_lower for k in ["队列", "queue"]):
        return TOPIC_OUTLINES["队列"]
    if any(k in topic_lower for k in ["图", "graph", "dfs", "bfs"]):
        return TOPIC_OUTLINES["图"]

    # 通用深度大纲
    return {
        "title": f"{topic} 详解",
        "subtitle": f"{subject} · 核心知识点",
        "slides": [
            {"type": "title", "title": f"{topic} 详解", "subtitle": f"{subject} · 从概念到实战"},
            {"type": "definition", "title": f"什么是{topic}", "concept": topic,
             "definition": f"{topic}是{subject}中的核心概念，掌握它是理解后续高级内容的基础。",
             "key_properties": [f"{topic}的基本特征", f"{topic}的核心性质", f"{topic}的分类方式"],
             "analogy": f"可以用生活中的例子来类比理解{topic}"},
            {"type": "structure", "title": f"{topic}的结构定义", "description": f"{topic}在C语言中的表示方式",
             "fields": [{"name": "data", "type": "int", "desc": "数据域"}, {"name": "next/link", "type": "指针", "desc": "链接域"}],
             "code": "// " + topic + "的C语言定义\ntypedef struct {\n    int data;\n    // 其他字段\n} " + topic.replace(' ', '') + "Type;"},
            {"type": "algorithm", "title": f"{topic}的基本操作", "description": f"{topic}最常用的操作及实现",
             "steps": [f"步骤1：初始化{topic}", "步骤2：执行核心操作", "步骤3：处理边界条件", "步骤4：返回结果"],
             "code": "// " + topic + "核心操作\nvoid operation(" + topic.replace(' ', '') + "Type *S) {\n    if (S == NULL) return;\n    // 核心逻辑\n    printf(\"操作完成\");\n}",
             "complexity": {"time": "O(n)", "space": "O(1)", "explanation": "需要遍历所有元素"}},
            {"type": "code_walkthrough", "title": f"{topic}关键代码解析",
             "code": "// " + topic + "的关键代码片段\nint key_code(" + topic.replace(' ', '') + "Type *S) {\n    int result = 0;\n    while (S != NULL) {\n        result += S->data;\n        S = S->next;\n    }\n    return result;\n}",
             "explanations": [{"lines": "第3行", "explain": "初始化结果变量"}, {"lines": "第4-7行", "explain": "循环遍历每个节点，累加数据"}, {"lines": "第8行", "explain": "返回最终结果"}]},
            {"type": "comparison", "title": f"{topic}与其他方案对比",
             "headers": ["对比项", f"{topic}", "替代方案"],
             "rows": [["时间复杂度", "O(n)", "O(nlogn)"], ["空间复杂度", "O(1)", "O(n)"], ["实现难度", "中等", "较高"], ["适用场景", "通用", "特定场景"]]},
            {"type": "knowledge_map", "title": f"{topic}知识结构", "items": [topic, "基本概念", "核心操作", "应用场景", "常见错误"],
             "relations": [[topic, "基本概念"], [topic, "核心操作"], ["核心操作", "应用场景"], ["基本概念", "常见错误"]]},
            {"type": "mistakes", "title": "常见易错点",
             "items": [
                 {"point": f"混淆{topic}的相关概念", "correct": "对比学习，明确每个概念的定义和区别"},
                 {"point": "忽略边界条件（空表、单元素等）", "correct": "写代码前先画图分析边界情况"},
                 {"point": f"算法步骤顺序错误", "correct": "按步骤逐一检查，特别是指针操作的顺序"}
             ]},
            {"type": "quiz", "title": "练习题",
             "questions": [
                 {"q": f"关于{topic}，以下说法正确的是？", "options": [f"A. {topic}不需要初始化", f"B. {topic}的核心是递归思想", f"C. {topic}的时间复杂度一定是O(1)", f"D. {topic}只能用数组实现"], "answer": "B",
                  "explain": f"{topic}的核心在于理解其递归/迭代的本质，其他选项都过于绝对"},
                 {"q": f"以下关于{topic}的时间复杂度，正确的是？", "options": ["A. 所有操作都是O(1)", "B. 查找操作是O(n)", "C. 插入操作是O(n²)", "D. 删除操作是O(1)"], "answer": "B",
                  "explain": f"{topic}的查找需要遍历，时间复杂度为O(n)"},
                 {"q": f"实现{topic}时最容易犯的错误是？", "options": ["A. 忘记初始化", "B. 指针操作顺序错误", "C. 没有处理空值", "D. 以上都是"], "answer": "D", "explain": "这三个都是常见错误，需要特别注意"}
             ]},
            {"type": "summary", "title": "总结", "points": [
                f"掌握了{topic}的核心概念和定义",
                f"学会了{topic}的基本操作实现",
                f"理解了{topic}的时间空间复杂度",
                "识别了常见错误和注意事项",
            ], "next_topic": "下一个相关主题"},
        ],
    }


# ============================================================
# PPTX 构建器 v2 — 丰富的幻灯片类型
# ============================================================

def build_pptx(outline: Dict[str, Any], output_path: str) -> str:
    """根据大纲生成 .pptx 文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 颜色主题
    C_PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
    C_PRIMARY_LIGHT = RGBColor(0x81, 0x8C, 0xF8)
    C_SECONDARY = RGBColor(0x0E, 0xA5, 0xE9)
    C_SUCCESS = RGBColor(0x10, 0xB9, 0x81)
    C_WARNING = RGBColor(0xF5, 0x9E, 0x0B)
    C_DANGER = RGBColor(0xEF, 0x44, 0x44)
    C_DARK = RGBColor(0x1E, 0x29, 0x3B)
    C_MEDIUM = RGBColor(0x47, 0x55, 0x69)
    C_LIGHT = RGBColor(0x64, 0x74, 0x8B)
    C_MUTED = RGBColor(0x94, 0xA3, 0xB8)
    C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    C_BG = RGBColor(0xF8, 0xFA, 0xFC)
    C_BG_CODE = RGBColor(0xF1, 0xF5, 0xF9)
    C_BORDER = RGBColor(0xE2, 0xE8, 0xF0)

    def add_bg(slide, color=C_BG):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_shape(slide, left, top, width, height, fill_color, border_color=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        if border_color:
            shape.line.fill.solid()
            shape.line.fill.fore_color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    def add_text(slide, left, top, width, height, text, size=18, color=C_DARK, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = align
        return txBox

    def add_multiline_text(slide, left, top, width, height, lines, size=16, color=C_DARK, line_spacing=1.3, font_name="Consolas"):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, line in enumerate(lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.name = font_name
            p.space_after = Pt(size * (line_spacing - 1) * 2)
        return txBox

    def add_code_block(slide, left, top, width, height, code, font_size=13):
        bg = add_shape(slide, left, top, width, height, C_BG_CODE, C_BORDER)
        lines = code.strip().split("\n")
        add_multiline_text(slide, left + Inches(0.2), top + Inches(0.15),
                          width - Inches(0.4), height - Inches(0.3),
                          lines, size=font_size, color=C_DARK, font_name="Consolas")

    def add_accent_bar(slide, left, top, height=Inches(0.5)):
        shape = slide.shapes.add_shape(1, left, top, Inches(0.07), height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = C_PRIMARY
        shape.line.fill.background()

    def add_page_number(slide, num, total):
        add_text(slide, Inches(12.2), Inches(7.0), Inches(1), Inches(0.4),
                f"{num}/{total}", size=11, color=C_MUTED, align=PP_ALIGN.RIGHT)

    def add_section_header(slide, title, subtitle=None):
        add_accent_bar(slide, Inches(0.8), Inches(0.55))
        add_text(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.7),
                title, size=30, color=C_DARK, bold=True)
        if subtitle:
            add_text(slide, Inches(1.1), Inches(1.05), Inches(10), Inches(0.4),
                    subtitle, size=14, color=C_LIGHT)

    total_slides = len(outline.get("slides", []))

    for idx, slide_data in enumerate(outline.get("slides", []), 1):
        slide_type = slide_data.get("type", "content")
        layout = prs.slide_layouts[6]  # blank

        # ── 标题页 ──
        if slide_type == "title":
            slide = prs.slides.add_slide(layout)
            add_bg(slide, C_PRIMARY)
            # 装饰
            add_shape(slide, Inches(9.5), Inches(-1.5), Inches(5.5), Inches(5.5),
                     C_PRIMARY_LIGHT, shape_type=MSO_SHAPE.OVAL)
            add_shape(slide, Inches(-1.5), Inches(4.5), Inches(4.5), Inches(4.5),
                     RGBColor(0x43, 0x38, 0xCA), shape_type=MSO_SHAPE.OVAL)
            # 标题
            add_text(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                    slide_data.get("title", ""), size=44, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
            add_text(slide, Inches(1.5), Inches(3.6), Inches(10), Inches(0.8),
                    slide_data.get("subtitle", ""), size=20, color=RGBColor(0xC7, 0xD2, 0xFE), align=PP_ALIGN.CENTER)
            add_text(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
                    "LearnLab · 智能生成", size=14, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER)

        # ── 概念定义页 ──
        elif slide_type == "definition":
            slide = prs.slides.add_slide(layout)
            add_bg(slide)
            add_section_header(slide, slide_data.get("title", "概念定义"))

            concept = slide_data.get("concept", "")
            definition = slide_data.get("definition", "")
            properties = slide_data.get("key_properties", [])
            analogy = slide_data.get("analogy", "")

            # 概念名 + 定义
            add_shape(slide, Inches(0.8), Inches(1.5), Inches(11.5), Inches(1.6), C_WHITE, C_BORDER)
            add_text(slide, Inches(1.1), Inches(1.6), Inches(3), Inches(0.5),
                    concept, size=22, color=C_PRIMARY, bold=True)
            add_text(slide, Inches(1.1), Inches(2.15), Inches(10.8), Inches(0.8),
                    definition, size=15, color=C_MEDIUM)

            # 性质卡片
            if properties:
                add_text(slide, Inches(0.8), Inches(3.4), Inches(3), Inches(0.4),
                        "核心性质", size=16, color=C_DARK, bold=True)
                card_w = Inches(min(3.6, 11.0 / max(len(properties), 1)))
                for i, prop in enumerate(properties[:4]):
                    x = Inches(0.8 + i * 3.9)
                    add_shape(slide, x, Inches(3.9), card_w, Inches(1.2), C_WHITE, C_BORDER)
                    num = add_shape(slide, x + Inches(0.15), Inches(4.05), Inches(0.4), Inches(0.4),
                                   C_PRIMARY, shape_type=MSO_SHAPE.OVAL)
                    tf = num.text_frame
                    tf.paragraphs[0].text = str(i + 1)
                    tf.paragraphs[0].font.size = Pt(14)
                    tf.paragraphs[0].font.color.rgb = C_WHITE
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    add_text(slide, x + Inches(0.65), Inches(4.0), card_w - Inches(0.8), Inches(1.0),
                            prop, size=14, color=C_MEDIUM)

            # 类比
            if analogy:
                add_shape(slide, Inches(0.8), Inches(5.5), Inches(11.5), Inches(1.0),
                         RGBColor(0xEF, 0xF6, 0xFF), C_SECONDARY)
                add_text(slide, Inches(1.1), Inches(5.55), Inches(1), Inches(0.4),
                        "💡 类比", size=14, color=C_SECONDARY, bold=True)
                add_text(slide, Inches(1.1), Inches(5.95), Inches(10.8), Inches(0.5),
                        analogy, size=15, color=C_MEDIUM)
            add_page_number(slide, idx, total_slides)

        # ── 数据结构/语法详解页 ──
        elif slide_type == "structure":
            slide = prs.slides.add_slide(layout)
            add_bg(slide)
            add_section_header(slide, slide_data.get("title", "结构详解"), slide_data.get("description"))

            fields = slide_data.get("fields", [])
            code = slide_data.get("code", "")

            # 字段说明表
            if fields:
                y_start = 1.6
                # 表头
                add_shape(slide, Inches(0.8), Inches(y_start), Inches(5.5), Inches(0.45), C_PRIMARY)
                for j, header in enumerate(["字段名", "类型", "说明"]):
                    x = Inches(0.9 + j * 1.8)
                    add_text(slide, x, Inches(y_start + 0.05), Inches(1.7), Inches(0.35),
                            header, size=13, color=C_WHITE, bold=True)
                for i, field in enumerate(fields):
                    y = y_start + 0.5 + i * 0.5
                    bg_color = C_WHITE if i % 2 == 0 else C_BG
                    add_shape(slide, Inches(0.8), Inches(y), Inches(5.5), Inches(0.45), bg_color)
                    add_text(slide, Inches(0.9), Inches(y + 0.05), Inches(1.6), Inches(0.35),
                            field.get("name", ""), size=13, color=C_PRIMARY, bold=True, font_name="Consolas")
                    add_text(slide, Inches(2.7), Inches(y + 0.05), Inches(1.6), Inches(0.35),
                            field.get("type", ""), size=13, color=C_SECONDARY, font_name="Consolas")
                    add_text(slide, Inches(4.5), Inches(y + 0.05), Inches(1.7), Inches(0.35),
                            field.get("desc", ""), size=13, color=C_MEDIUM)

            # 代码
            if code:
                add_text(slide, Inches(7), Inches(1.6), Inches(5), Inches(0.4),
                        "C语言定义", size=14, color=C_DARK, bold=True)
                add_code_block(slide, Inches(7), Inches(2.1), Inches(5.5), Inches(4.5), code, font_size=14)
            add_page_number(slide, idx, total_slides)

        # ── 算法步骤页 ──
        elif slide_type == "algorithm":
            slide = prs.slides.add_slide(layout)
            add_bg(slide)
            add_section_header(slide, slide_data.get("title", "算法"), slide_data.get("description"))

            steps = slide_data.get("steps", [])
            code = slide_data.get("code", "")
            complexity = slide_data.get("complexity", {})

            # 步骤
            if steps:
                add_text(slide, Inches(0.8), Inches(1.5), Inches(2), Inches(0.35),
                        "执行步骤", size=14, color=C_DARK, bold=True)
                for i, step in enumerate(steps[:6]):
                    y = 1.9 + i * 0.55
                    add_shape(slide, Inches(0.8), Inches(y), Inches(5.5), Inches(0.45), C_WHITE, C_BORDER)
                    num = add_shape(slide, Inches(0.9), Inches(y + 0.05), Inches(0.32), Inches(0.32),
                                   C_PRIMARY, shape_type=MSO_SHAPE.OVAL)
                    tf = num.text_frame
                    tf.paragraphs[0].text = str(i + 1)
                    tf.paragraphs[0].font.size = Pt(11)
                    tf.paragraphs[0].font.color.rgb = C_WHITE
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    add_text(slide, Inches(1.35), Inches(y + 0.05), Inches(4.8), Inches(0.35),
                            step, size=13, color=C_MEDIUM)

            # 代码
            if code:
                add_text(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.35),
                        "C语言实现", size=14, color=C_DARK, bold=True)
                code_h = Inches(4.0) if not complexity else Inches(3.2)
                add_code_block(slide, Inches(7), Inches(1.9), Inches(5.5), code_h, code, font_size=13)

            # 复杂度
            if complexity:
                cy = 5.3 if code else 1.9
                add_shape(slide, Inches(7), Inches(cy), Inches(5.5), Inches(1.3),
                         RGBColor(0xFF, 0xFB, 0xEB), C_WARNING)
                add_text(slide, Inches(7.2), Inches(cy + 0.05), Inches(2), Inches(0.35),
                        "复杂度分析", size=13, color=C_WARNING, bold=True)
                time_str = complexity.get("time", "O(n)")
                space_str = complexity.get("space", "O(1)")
                add_text(slide, Inches(7.2), Inches(cy + 0.4), Inches(2), Inches(0.3),
                        f"时间: {time_str}", size=14, color=C_DARK, bold=True, font_name="Consolas")
                add_text(slide, Inches(9.5), Inches(cy + 0.4), Inches(2.5), Inches(0.3),
                        f"空间: {space_str}", size=14, color=C_DARK, bold=True, font_name="Consolas")
                expl = complexity.get("explanation", "")
                if expl:
                    add_text(slide, Inches(7.2), Inches(cy + 0.75), Inches(5.1), Inches(0.5),
                            expl, size=12, color=C_LIGHT)
            add_page_number(slide, idx, total_slides)

        # ── 代码逐行解析页 ──
        elif slide_type == "code_walkthrough":
            slide = prs.slides.add_slide(layout)
            add_bg(slide)
            add_section_header(slide, slide_data.get("title", "代码解析"))

            code = slide_data.get("code", "")
            explanations = slide_data.get("explanations", [])

            if code:
                add_code_block(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), code, font_size=13)

            if explanations:
                add_text(slide, Inches(7), Inches(1.5), Inches(5), Inches(0.35),
                        "逐行解析", size=14, color=C_DARK, bold=True)
                for i, exp in enumerate(explanations[:6]):
                    y = 2.0 + i * 0.85
                    add_shape(slide, Inches(7), Inches(y), Inches(5.5), Inches(0.7), C_WHITE, C_BORDER)
                    add_text(slide, Inches(7.2), Inches(y + 0.05), Inches(5), Inches(0.3),
                            exp.get("lines", ""), size=12, color=C_PRIMARY, bold=True, font_name="Consolas")
                    add_text(slide, Inches(7.2), Inches(y + 0.32), Inches(5), Inches(0.35),
                            exp.get("explain", ""), size=12, color=C_MEDIUM)
            add_page_number(slide, idx, total_slides)

        # ── 对比分析页 ──
        elif slide_type == "comparison":
            slide = prs.slides.add_slide(layout)
            add_bg(slide)
            add_section_header(slide, slide_data.get("title", "对比分析"))

            headers = slide_data.get("headers", [])
            rows = slide_data.get("rows", [])

            if headers and rows:
                col_count = len(headers)
                col_w = min(3.0, 11.0 / col_count)
                x_start = 0.8
                # 表头
                add_shape(slide, Inches(x_start), Inches(1.5), Inches(col_w * col_count), Inches(0.5), C_PRIMARY)
                for j, h in enumerate(headers):
                    add_text(slide, Inches(x_start + j * col_w + 0.1), Inches(1.55), Inches(col_w - 0.2), Inches(0.4),
                            h, size=14, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
                # 数据行
                for i, row in enumerate(rows[:8]):
                    y = 2.05 + i * 0.55
                    bg = C_WHITE if i % 2 == 0 else C_BG
                    add_shape(slide, Inches(x_start), Inches(y), Inches(col_w * col_count), Inches(0.5), bg, C_BORDER)
                    for j, cell in enumerate(row[:col_count]):
                        is_first = j == 0
                        add_text(slide, Inches(x_start + j * col_w + 0.1), Inches(y + 0.05),
                                Inches(col_w - 0.2), Inches(0.4),
                                str(cell), size=13,
                                color=C_DARK if is_first else C_MEDIUM,
                                bold=is_first,
                                align=PP_ALIGN.CENTER)
            add_page_number(slide, idx, total_slides)

        # ── 知识结构图 ──
        elif slide_type == "knowledge_map":
            slide = prs.slides.add_slide(layout)
            add_bg(slide)
            add_section_header(slide, slide_data.get("title", "知识结构图"))

            items = slide_data.get("items", [])
            if items:
                import math
                cx, cy = Inches(6.5), Inches(4.0)
                # 中心节点
                center = add_shape(slide, cx - Inches(1.2), cy - Inches(0.4), Inches(2.4), Inches(0.85), C_PRIMARY)
                tf = center.text_frame
                tf.paragraphs[0].text = items[0]
                tf.paragraphs[0].font.size = Pt(18)
                tf.paragraphs[0].font.color.rgb = C_WHITE
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER

                node_colors = [C_SECONDARY, C_SUCCESS, C_WARNING, C_DANGER, C_PRIMARY_LIGHT]
                for i, item in enumerate(items[1:]):
                    angle = (2 * math.pi * i) / max(len(items) - 1, 1) - math.pi / 2
                    r = Inches(2.8)
                    nx = cx + int(r * math.cos(angle))
                    ny = cy + int(r * math.sin(angle))
                    node = add_shape(slide, nx - Inches(0.9), ny - Inches(0.3), Inches(1.8), Inches(0.65),
                                    node_colors[i % len(node_colors)])
                    tf = node.text_frame
                    tf.paragraphs[0].text = item
                    tf.paragraphs[0].font.size = Pt(14)
                    tf.paragraphs[0].font.color.rgb = C_WHITE
                    tf.paragraphs[0].font.bold = True
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            add_page_number(slide, idx, total_slides)

        # ── 易错点页 ──
        elif slide_type == "mistakes":
            slide = prs.slides.add_slide(layout)
            add_bg(slide)
            add_section_header(slide, slide_data.get("title", "常见易错点"))

            items = slide_data.get("items", [])
            for i, item in enumerate(items[:4]):
                y = 1.6 + i * 1.4
                h = 1.2 if not item.get("code_example") else 1.3
                # 错误卡片
                add_shape(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(h), C_WHITE, RGBColor(0xFE, 0xE2, 0xE2))
                add_text(slide, Inches(1.1), Inches(y + 0.08), Inches(10.5), Inches(0.35),
                        f"❌  {item.get('point', '')}", size=15, color=C_DANGER, bold=True)
                add_text(slide, Inches(1.1), Inches(y + 0.42), Inches(10.5), Inches(0.35),
                        f"✅  {item.get('correct', '')}", size=14, color=C_SUCCESS)
                # 代码示例
                code_ex = item.get("code_example", "")
                if code_ex and y + h + 0.8 < 7.0:
                    add_code_block(slide, Inches(1.1), Inches(y + 0.8), Inches(10.8), Inches(0.8), code_ex, font_size=11)
            add_page_number(slide, idx, total_slides)

        # ── 练习题页 ──
        elif slide_type == "quiz":
            slide = prs.slides.add_slide(layout)
            add_bg(slide)
            add_section_header(slide, slide_data.get("title", "练习题"))

            questions = slide_data.get("questions", [])
            for i, q in enumerate(questions[:3]):
                y = 1.5 + i * 1.9
                # 题号
                badge = add_shape(slide, Inches(0.8), Inches(y), Inches(0.5), Inches(0.4), C_PRIMARY, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE)
                tf = badge.text_frame
                tf.paragraphs[0].text = f"Q{i+1}"
                tf.paragraphs[0].font.size = Pt(12)
                tf.paragraphs[0].font.color.rgb = C_WHITE
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                # 题目
                add_text(slide, Inches(1.4), Inches(y), Inches(11), Inches(0.4),
                        q.get("q", ""), size=15, color=C_DARK, bold=True)
                # 选项
                options = q.get("options", [])
                for j, opt in enumerate(options[:4]):
                    add_text(slide, Inches(1.6), Inches(y + 0.4 + j * 0.3), Inches(10), Inches(0.3),
                            opt, size=13, color=C_MEDIUM)
                # 答案
                ans_y = y + 0.4 + min(len(options), 4) * 0.3 + 0.1
                if q.get("answer"):
                    add_shape(slide, Inches(1.4), Inches(ans_y), Inches(10.5), Inches(0.5),
                             RGBColor(0xEF, 0xF6, 0xFF), C_SECONDARY)
                    add_text(slide, Inches(1.6), Inches(ans_y + 0.05), Inches(10), Inches(0.4),
                            f"答案: {q['answer']}  |  {q.get('explain', '')}", size=12, color=C_PRIMARY)
            add_page_number(slide, idx, total_slides)

        # ── 总结页 ──
        elif slide_type == "summary":
            slide = prs.slides.add_slide(layout)
            add_bg(slide, C_PRIMARY)
            add_shape(slide, Inches(10), Inches(4.5), Inches(5), Inches(5),
                     C_PRIMARY_LIGHT, shape_type=MSO_SHAPE.OVAL)
            add_text(slide, Inches(1.5), Inches(0.8), Inches(10), Inches(1),
                    slide_data.get("title", "总结"), size=40, color=C_WHITE, bold=True, align=PP_ALIGN.CENTER)
            points = slide_data.get("points", [])
            for i, point in enumerate(points[:6]):
                y = 2.2 + i * 0.85
                add_shape(slide, Inches(1.5), Inches(y), Inches(10), Inches(0.65),
                         RGBColor(0x43, 0x38, 0xCA))
                add_text(slide, Inches(1.8), Inches(y + 0.1), Inches(9.5), Inches(0.45),
                        f"✦  {point}", size=20, color=RGBColor(0xE0, 0xE7, 0xFF))
            next_topic = slide_data.get("next_topic", "")
            if next_topic:
                add_text(slide, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
                        f"下一讲: {next_topic}", size=16, color=RGBColor(0xA5, 0xB4, 0xFC), align=PP_ALIGN.CENTER)
            add_page_number(slide, idx, total_slides)

        # ── 兼容旧类型 ──
        elif slide_type == "content":
            slide = prs.slides.add_slide(layout)
            add_bg(slide)
            add_section_header(slide, slide_data.get("title", ""))
            bullets = slide_data.get("bullets", [])
            for i, bullet in enumerate(bullets[:6]):
                y = 1.6 + i * 0.9
                add_shape(slide, Inches(0.8), Inches(y), Inches(11.5), Inches(0.7), C_WHITE, C_BORDER)
                num = add_shape(slide, Inches(1.0), Inches(y + 0.12), Inches(0.42), Inches(0.42),
                               C_PRIMARY, shape_type=MSO_SHAPE.OVAL)
                tf = num.text_frame
                tf.paragraphs[0].text = str(i + 1)
                tf.paragraphs[0].font.size = Pt(14)
                tf.paragraphs[0].font.color.rgb = C_WHITE
                tf.paragraphs[0].font.bold = True
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                add_text(slide, Inches(1.6), Inches(y + 0.12), Inches(10), Inches(0.45),
                        bullet, size=16, color=C_DARK)
            add_page_number(slide, idx, total_slides)

    prs.save(output_path)
    return output_path


# ============================================================
# 完整生成流程
# ============================================================

async def generate_ppt(
    topic: str,
    subject: str = "C语言数据结构",
    llm=None,
) -> Dict[str, Any]:
    """完整 PPT 生成流程"""
    if llm is None:
        try:
            from ..services.llm_factory import LLMFactory
            llm = LLMFactory.get_default_llm()
        except Exception as e:
            logger.warning(f"PPT生成LLM初始化失败，将使用默认大纲: {e}")

    outline = await generate_ppt_outline(topic, subject, llm)

    timestamp = datetime.now(timezone.utc).strftime("%Y%m%d_%H%M%S")
    safe_topic = re.sub(r'[^\w一-鿿]', '_', topic)[:20]
    filename = f"{safe_topic}_{timestamp}.pptx"
    output_path = os.path.join(PPT_OUTPUT_DIR, filename)

    build_pptx(outline, output_path)

    return {
        "filename": filename,
        "path": output_path,
        "outline": outline,
        "slide_count": len(outline.get("slides", [])),
    }
